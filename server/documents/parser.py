import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pypdf import PdfReader


TEXT_LIKE_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".json",
    ".html",
    ".css",
    ".csv",
    ".yml",
    ".yaml",
    ".toml",
    ".xml",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".cs",
    ".go",
    ".rs",
    ".sql",
    ".sh",
    ".ps1",
}


def decode_text(raw: bytes) -> str:
    for encoding in ("utf-8", "gbk"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def normalize_text(text: str) -> str:
    return text.replace("\r\n", "\n").replace("\r", "\n").strip()


def split_paragraphs(text: str) -> List[str]:
    parts = re.split(r"\n\s*\n+", text)
    return [part.strip() for part in parts if part.strip()]


def chunk_text(
    text: str,
    chunk_size: int = 600,
    overlap: int = 80,
    page_no: Optional[int] = None,
    start_index: int = 0,
    min_chunk_chars: int = 20,
) -> List[Dict[str, Any]]:
    if chunk_size <= 0:
        raise ValueError("chunk_size must be greater than 0")
    if overlap < 0 or overlap >= chunk_size:
        raise ValueError("overlap must be >= 0 and < chunk_size")

    text = normalize_text(text)
    if not text:
        return []

    paragraphs = split_paragraphs(text) or [text]
    chunks: List[Dict[str, Any]] = []
    idx = start_index
    step = chunk_size - overlap

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(para) <= chunk_size:
            if len(para) >= min_chunk_chars:
                chunks.append(
                    {
                        "chunk_index": idx,
                        "content": para,
                        "token_count": len(para),
                        "page_no": page_no,
                        "tags": None,
                    }
                )
                idx += 1
            continue

        start = 0
        while start < len(para):
            piece = para[start : start + chunk_size].strip()
            if len(piece) >= min_chunk_chars:
                chunks.append(
                    {
                        "chunk_index": idx,
                        "content": piece,
                        "token_count": len(piece),
                        "page_no": page_no,
                        "tags": None,
                    }
                )
                idx += 1
            start += step

    return chunks


def parse_txt(
    file_path: str,
    chunk_size: int = 600,
    overlap: int = 80,
) -> List[Dict[str, Any]]:
    raw = Path(file_path).read_bytes()
    text = decode_text(raw)
    chunks = chunk_text(
        text,
        chunk_size=chunk_size,
        overlap=overlap,
        page_no=None,
        start_index=0,
    )
    if not chunks:
        raise ValueError("txt 文件为空或无可解析文本")
    return chunks


def parse_text_like(
    file_path: str,
    chunk_size: int = 600,
    overlap: int = 80,
) -> List[Dict[str, Any]]:
    return parse_txt(file_path, chunk_size=chunk_size, overlap=overlap)


def parse_pdf(
    file_path: str,
    chunk_size: int = 600,
    overlap: int = 80,
) -> List[Dict[str, Any]]:
    reader = PdfReader(str(Path(file_path)))
    chunks: List[Dict[str, Any]] = []
    next_index = 0

    for page_no, page in enumerate(reader.pages, start=1):
        page_text = normalize_text(page.extract_text() or "")
        if not page_text:
            continue

        page_chunks = chunk_text(
            page_text,
            chunk_size=chunk_size,
            overlap=overlap,
            page_no=page_no,
            start_index=next_index,
        )
        chunks.extend(page_chunks)
        next_index += len(page_chunks)

    if not chunks:
        raise ValueError("pdf 文件为空或无法提取文本")
    return chunks


def parse_docx(
    file_path: str,
    chunk_size: int = 600,
    overlap: int = 80,
) -> List[Dict[str, Any]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise ValueError("python-docx is not installed") from exc

    doc = Document(str(Path(file_path)))
    lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                lines.append(" | ".join(cells))

    chunks = chunk_text(
        "\n\n".join(lines),
        chunk_size=chunk_size,
        overlap=overlap,
        page_no=None,
        start_index=0,
    )
    if not chunks:
        raise ValueError("docx 文件为空或无可解析文本")
    return chunks


def parse_pptx(
    file_path: str,
    chunk_size: int = 600,
    overlap: int = 80,
) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("python-pptx is not installed") from exc

    prs = Presentation(str(Path(file_path)))
    chunks: List[Dict[str, Any]] = []
    next_index = 0
    for slide_no, slide in enumerate(prs.slides, start=1):
        lines: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
        if not lines:
            continue
        slide_chunks = chunk_text(
            "\n\n".join(lines),
            chunk_size=chunk_size,
            overlap=overlap,
            page_no=slide_no,
            start_index=next_index,
        )
        chunks.extend(slide_chunks)
        next_index += len(slide_chunks)

    if not chunks:
        raise ValueError("pptx 文件为空或无可解析文本")
    return chunks


def parse_document(
    file_path: str,
    chunk_size: int = 600,
    overlap: int = 80,
) -> Tuple[str, List[Dict[str, Any]]]:
    suffix = Path(file_path).suffix.lower()

    if suffix in TEXT_LIKE_SUFFIXES:
        return suffix.lstrip("."), parse_text_like(file_path, chunk_size=chunk_size, overlap=overlap)
    if suffix == ".pdf":
        return "pdf", parse_pdf(file_path, chunk_size=chunk_size, overlap=overlap)
    if suffix == ".docx":
        return "docx", parse_docx(file_path, chunk_size=chunk_size, overlap=overlap)
    if suffix == ".pptx":
        return "pptx", parse_pptx(file_path, chunk_size=chunk_size, overlap=overlap)

    raise ValueError(f"暂不支持的文件类型: {suffix}")
